import streamlit as st
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import io

# Streamlit App Title
st.title("Interactive PowerPoint Generator")
st.write("Customize your PowerPoint presentation by adjusting sections, colors, and slides per section.")

# Section Defaults
default_sections = [
    {"title": "Index", "slides": 1, "color": "#0070C0"},
    {"title": "Introduction", "slides": 2, "color": "#7030A0"},
    {"title": "PEIR Scope", "slides": 3, "color": "#C00000"},
    {"title": "Deliverables", "slides": 2, "color": "#00B050"},
    {"title": "Qualifications", "slides": 2, "color": "#FFC000"},
    {"title": "Timeline & Budget", "slides": 1, "color": "#7F7F7F"},
    {"title": "Closing", "slides": 1, "color": "#0070C0"},
]

# User Input for Sections
st.sidebar.header("Customize Sections")
sections = []
for i, section in enumerate(default_sections):
    with st.sidebar.expander(f"Section {i + 1}: {section['title']}"):
        title = st.text_input(f"Title for Section {i + 1}", section["title"])
        slides = st.number_input(f"Number of slides for {title}", min_value=1, max_value=10, value=section["slides"])
        color = st.color_picker(f"Pick a color for {title}", section["color"])
        sections.append({"title": title, "slides": slides, "color": color})

# Function to convert hex color to RGBColor
def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip("#")
    return RGBColor(*(int(hex_color[i : i + 2], 16) for i in (0, 2, 4)))

# Function to add a section divider
def add_section_divider(prs, section_title, color):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(color)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(7.5), Inches(2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.add_paragraph()
    p.text = section_title
    p.font.name = "Helvetica"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

# Function to add content slides
def add_slide_with_template(prs, title, subtitle, details, section_index, prompts, include_page_number=True, slide_number=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    rectangle = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left=Inches(6.5), top=Inches(0.3), width=Inches(1.5), height=Inches(0.1)
    )
    rectangle.fill.solid()
    rectangle.fill.fore_color.rgb = hex_to_rgb(sections[section_index]["color"])
    rectangle.line.color.rgb = RGBColor(255, 255, 255)
    rectangle.shadow.inherit = False

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(7.5), Inches(1))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.add_paragraph()
    p.text = title
    p.font.name = "Helvetica"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(7.5), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.add_paragraph()
    p.text = subtitle
    p.font.name = "Helvetica"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(50, 50, 50)

    details_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(7.5), Inches(3))
    details_frame = details_box.text_frame
    details_frame.word_wrap = True
    p = details_frame.add_paragraph()
    p.text = details
    p.font.name = "Helvetica"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(150, 150, 150)

    prompt_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(7.5), Inches(1))
    prompt_frame = prompt_box.text_frame
    p = prompt_frame.add_paragraph()
    p.text = "Consider the Following:"
    p.font.name = "Helvetica"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    for prompt in prompts:
        p = prompt_frame.add_paragraph()
        p.text = f"• {prompt}"
        p.font.name = "Helvetica"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(50, 50, 50)

    if include_page_number and slide_number is not None:
        page_number_box = slide.shapes.add_textbox(Inches(7.5), Inches(10.5), Inches(0.5), Inches(0.5))
        page_number_frame = page_number_box.text_frame
        page_number_frame.word_wrap = False
        page_number_frame.margin_left = 0
        page_number_frame.margin_right = 0
        page_number_frame.margin_top = 0
        page_number_frame.margin_bottom = 0
        p = page_number_frame.add_paragraph()
        p.text = str(slide_number)
        p.font.name = "Helvetica"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = 2

# Generate PowerPoint Button
if st.button("Generate PowerPoint"):
    prs = Presentation()
    prs.slide_width = Inches(8.5)
    prs.slide_height = Inches(11)

    # Add title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(7.5), Inches(2))
    title_frame = title_box.text_frame
    p = title_frame.add_paragraph()
    p.text = "Customized PowerPoint Presentation"
    p.font.name = "Helvetica"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(7.5), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.add_paragraph()
    p.text = "Generated by Streamlit App"
    p.font.name = "Helvetica"
    p.font.size = Pt(20)
    p.font.bold = False
    p.font.color.rgb = RGBColor(50, 50, 50)

    # Add sections dynamically
    slide_number = 2
    for i, section in enumerate(sections):
        add_section_divider(prs, section["title"], section["color"])
        for j in range(section["slides"]):
            add_slide_with_template(
                prs,
                f"{section['title']} - Slide {j + 1}",
                f"Subtitle for {section['title']}",
                f"Details for {section['title']} slide {j + 1}.",
                i,
                [f"Prompt {j + 1} for {section['title']}."],
                slide_number=slide_number,
            )
            slide_number += 1

    # Save to buffer and offer download
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    st.download_button("Download PowerPoint", data=buffer, file_name="Customized_Presentation.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
